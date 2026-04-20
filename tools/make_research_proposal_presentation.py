#!/usr/bin/env python3
"""Build the preliminary research proposal presentation.

Run:
    conda activate svi_segformer
    python tools/make_research_proposal_presentation.py
"""

from __future__ import annotations

import csv
import json
import math
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
import pandas as pd
from PIL import Image, ImageDraw, ImageEnhance, ImageFilter, ImageOps


ROOT = Path(__file__).resolve().parents[1]
FIG_DIR = ROOT / "outputs" / "figures" / "presentation"
PPT_DIR = ROOT / "papers" / "presentation"
FIG_DIR.mkdir(parents=True, exist_ok=True)
PPT_DIR.mkdir(parents=True, exist_ok=True)

OUT_PPT = PPT_DIR / "from_streetview_to_aerial_research_proposal.pptx"

NAVY = "#0B1F3A"
NAVY2 = "#002060"
BLUE = "#2E75B6"
LT_BLUE = "#D8E6F4"
GREEN = "#2F8F68"
LT_GREEN = "#DDEFE5"
ORANGE = "#F28E2B"
LT_ORANGE = "#FCE4D6"
RED = "#C00000"
DARK = "#262626"
GRAY = "#666666"
MID_GRAY = "#BFBFBF"
LIGHT = "#F5F7FA"
WHITE = "#FFFFFF"

SPLIT_INFO = {
    "p01_guro_digital": ("Val", ORANGE),
    "p02_songpa_jamsil": ("Test", RED),
    "p03_yeongdeungpo_mullae": ("Train", BLUE),
    "p04_seongdong_wangsimni": ("Test", RED),
    "p05_mapo_hongdae_yeonnam": ("Train", BLUE),
    "p06_gangnam_teheran": ("Train", BLUE),
}

URBAN_POOLS = [
    ("mapo_hongdae_yeonnam", 126.912, 37.548, 126.936, 37.566),
    ("seodaemun_ehwa_sinchon", 126.936, 37.555, 126.957, 37.566),
    ("jung_euljiro_chungmuro", 126.985, 37.557, 127.010, 37.570),
    ("seongdong_wangsimni", 127.030, 37.545, 127.055, 37.565),
    ("gangnam_teheran", 127.025, 37.497, 127.055, 37.510),
    ("songpa_jamsil", 127.075, 37.500, 127.105, 37.515),
    ("yeongdeungpo_mullae", 126.885, 37.505, 126.910, 37.525),
    ("guro_digital", 126.880, 37.475, 126.905, 37.495),
    ("dongdaemun_jangan", 127.055, 37.565, 127.080, 37.585),
]

META_PATH = ROOT / "data" / "raw" / "random_seoul_urban" / "dataset_metadata.json"
PATCH_CSV = ROOT / "data" / "raw" / "random_seoul_urban" / "sampling_patches.csv"
PATCH_TIF_DIR = ROOT / "data" / "raw" / "random_seoul_urban" / "patches"
SELECTED_CSV = ROOT / "data" / "processed" / "labeling_candidates_512_random_seoul_urban" / "selected_tiles.csv"
TILE_IMG_DIR = ROOT / "data" / "processed" / "labeling_candidates_512_random_seoul_urban" / "images"
SEM_IMG_DIR = ROOT / "data" / "processed" / "semantic_seg_512_random_seoul_urban" / "images"
SEM_MASK_DIR = ROOT / "data" / "processed" / "semantic_seg_512_random_seoul_urban" / "masks"
METRICS_CSV = ROOT / "outputs" / "segformer" / "tree_canopy_semantic_tcd_mit_b2_finetune_e50" / "metrics.csv"
THRESH_CSV = ROOT / "outputs" / "segformer" / "tree_canopy_semantic_tcd_mit_b2_finetune_e50" / "threshold_metrics.csv"
TCD_RUN_DIR = ROOT / "outputs" / "segformer" / "tree_canopy_semantic_tcd_mit_b2_finetune_e50"
PREVIEW_SHEET = (
    ROOT
    / "outputs"
    / "segformer"
    / "tree_canopy_semantic_tcd_mit_b2_finetune_e50"
    / "test_threshold_preview_sheet.jpg"
)


def read_csv_rows(path: Path) -> list[dict[str, str]]:
    with path.open(newline="", encoding="utf-8") as f:
        return list(csv.DictReader(f))


def read_patch_rows() -> list[dict[str, str]]:
    return read_csv_rows(PATCH_CSV)


def patch_short_name(pid: str) -> str:
    return pid.replace("p01_", "").replace("p02_", "").replace("p03_", "").replace("p04_", "").replace("p05_", "").replace("p06_", "").replace("_", " ")


def find_split_and_path(tile_stem: str) -> tuple[str | None, Path | None, Path | None]:
    for split in ("train", "val", "test"):
        img = SEM_IMG_DIR / split / f"{tile_stem}.png"
        mask = SEM_MASK_DIR / split / f"{tile_stem}.png"
        if img.exists():
            return split, img, mask if mask.exists() else None
    return None, None, None


def open_rgb(path: Path) -> Image.Image:
    img = Image.open(path).convert("RGB")
    return ImageEnhance.Contrast(img).enhance(1.08)


def read_patch_image(path: Path) -> Image.Image:
    try:
        import rasterio

        with rasterio.open(path) as src:
            arr = src.read([1, 2, 3])
        arr = np.moveaxis(arr, 0, -1).astype(np.float32)
        lo = np.percentile(arr, 2, axis=(0, 1))
        hi = np.percentile(arr, 98, axis=(0, 1))
        arr = np.clip((arr - lo) / np.maximum(hi - lo, 1), 0, 1)
        arr = (arr * 255).astype(np.uint8)
        return Image.fromarray(arr, "RGB")
    except Exception:
        return open_rgb(path)


def add_mask_overlay(img: Image.Image, mask: Image.Image, color=(242, 142, 43), alpha=0.48) -> Image.Image:
    base = img.convert("RGBA")
    m = mask.convert("L")
    overlay = Image.new("RGBA", base.size, (0, 0, 0, 0))
    rgba = (*color, int(255 * alpha))
    overlay.paste(rgba, mask=m.point(lambda v: 255 if v > 0 else 0))

    edge = m.filter(ImageFilter.FIND_EDGES).point(lambda v: 255 if v > 12 else 0)
    edge_layer = Image.new("RGBA", base.size, (*color, 255))
    base.alpha_composite(overlay)
    base.alpha_composite(Image.composite(edge_layer, Image.new("RGBA", base.size, (0, 0, 0, 0)), edge))
    return base.convert("RGB")


def save_fig(fig: plt.Figure, name: str, dpi: int = 170) -> Path:
    path = FIG_DIR / name
    fig.savefig(path, dpi=dpi, bbox_inches="tight", facecolor=WHITE)
    plt.close(fig)
    print(f"  {path.relative_to(ROOT)}")
    return path


def make_motivation_diagram() -> Path:
    fig, ax = plt.subplots(figsize=(12.5, 3.9), facecolor=WHITE)
    ax.set_xlim(0, 12.5)
    ax.set_ylim(0, 3.9)
    ax.axis("off")

    ax.text(
        6.25,
        3.65,
        "Do perceptual street-view metrics tell the same thermal story as physical aerial measurements?",
        ha="center",
        va="top",
        fontsize=14,
        color=NAVY,
        fontweight="bold",
    )

    panels = [
        (0.35, 0.35, 4.8, 2.75, "Street-view perception", "GVI / BVI / SVF", LT_BLUE, BLUE),
        (7.35, 0.35, 4.8, 2.75, "Aerial physical measurement", "Canopy / footprint / height", LT_GREEN, GREEN),
    ]
    for x, y, w, h, title, subtitle, fc, ec in panels:
        box = mpatches.FancyBboxPatch(
            (x, y),
            w,
            h,
            boxstyle="round,pad=0.05,rounding_size=0.06",
            facecolor=fc,
            edgecolor=ec,
            linewidth=1.8,
        )
        ax.add_patch(box)
        ax.text(x + 0.25, y + h - 0.38, title, fontsize=13, color=ec, fontweight="bold")
        ax.text(x + 0.25, y + 0.28, subtitle, fontsize=11, color=DARK, fontweight="bold")

    # Left side: eye / street canyon.
    ax.add_patch(mpatches.Ellipse((1.15, 1.8), 0.75, 0.38, facecolor=WHITE, edgecolor=BLUE, linewidth=1.5))
    ax.add_patch(mpatches.Circle((1.15, 1.8), 0.12, facecolor=BLUE, edgecolor=BLUE))
    for angle, color in [(12, GREEN), (0, GRAY), (-12, BLUE)]:
        ax.plot(
            [1.55, 4.65],
            [1.8, 1.8 + math.tan(math.radians(angle)) * 3.1],
            color=color,
            linewidth=2,
            alpha=0.8,
        )
    ax.add_patch(mpatches.Rectangle((3.95, 0.85), 0.35, 1.45, facecolor="#A7A9AC", edgecolor=GRAY))
    ax.add_patch(mpatches.Rectangle((4.4, 0.65), 0.5, 1.85, facecolor="#8C8C8C", edgecolor=GRAY))
    ax.add_patch(mpatches.Circle((3.45, 1.25), 0.34, facecolor="#78B66F", edgecolor=GREEN))
    ax.add_patch(mpatches.Rectangle((3.4, 0.62), 0.08, 0.65, facecolor="#7A4E2A", edgecolor="#7A4E2A"))
    ax.text(0.8, 0.8, "horizontal field of view", fontsize=9, color=GRAY, style="italic")

    # Right side: top-down canopy and buildings.
    rng = np.random.default_rng(4)
    for _ in range(18):
        x = rng.uniform(8.0, 11.55)
        y = rng.uniform(0.9, 2.65)
        r = rng.uniform(0.11, 0.22)
        ax.add_patch(mpatches.Circle((x, y), r, facecolor="#78B66F", edgecolor=GREEN, alpha=0.9))
    for x, y, w, h in [(8.1, 1.0, 0.45, 0.55), (9.0, 2.0, 0.6, 0.35), (10.2, 1.1, 0.55, 0.8), (11.0, 2.1, 0.65, 0.32)]:
        ax.add_patch(mpatches.Rectangle((x, y), w, h, facecolor="#B8B8B8", edgecolor=GRAY, linewidth=1))
    ax.text(8.0, 0.8, "top-down area / volume", fontsize=9, color=GRAY, style="italic")

    ax.annotate(
        "",
        xy=(7.0, 1.7),
        xytext=(5.45, 1.7),
        arrowprops=dict(arrowstyle="->", color=NAVY, linewidth=2),
    )
    ax.text(6.22, 2.0, "compare\nwithin same AT model", ha="center", fontsize=10, color=NAVY)
    return save_fig(fig, "proposal_motivation_diagram.png")


def make_pipeline_diagram() -> Path:
    steps = [
        ("V-World WMTS\n~0.25 m/px", LT_BLUE, BLUE),
        ("Random urban\npatch sampling", LT_BLUE, BLUE),
        ("512 px tiles\ncandidate selection", LT_BLUE, BLUE),
        ("Manual canopy\npolygon labeling", LT_GREEN, GREEN),
        ("Binary semantic\nmask conversion", LT_GREEN, GREEN),
        ("TCD SegFormer-B2\nfine-tuning", "#FFF2CC", "#9A6A00"),
        ("Tree Canopy Map\n0.25 m/px", "#FCE4D6", ORANGE),
        ("S-DoT 50 m\nbuffer join", "#EADCF8", "#7030A0"),
        ("Mixed-effects\nAT model", "#E2F0D9", GREEN),
    ]
    fig, ax = plt.subplots(figsize=(13.5, 2.6), facecolor=WHITE)
    ax.set_xlim(0, 13.5)
    ax.set_ylim(0, 2.6)
    ax.axis("off")
    box_w, box_h = 1.25, 1.28
    x0, y0, gap = 0.18, 0.65, 0.22
    for i, (text, fc, ec) in enumerate(steps):
        x = x0 + i * (box_w + gap)
        rect = mpatches.FancyBboxPatch(
            (x, y0),
            box_w,
            box_h,
            boxstyle="round,pad=0.04,rounding_size=0.05",
            facecolor=fc,
            edgecolor=ec,
            linewidth=1.5,
        )
        ax.add_patch(rect)
        ax.text(x + box_w / 2, y0 + box_h / 2, text, ha="center", va="center", fontsize=8.3, color=DARK, fontweight="bold")
        if i < len(steps) - 1:
            ax.annotate(
                "",
                xy=(x + box_w + gap * 0.8, y0 + box_h / 2),
                xytext=(x + box_w + gap * 0.15, y0 + box_h / 2),
                arrowprops=dict(arrowstyle="->", color=NAVY, linewidth=1.25),
            )
    ax.text(10.2, 2.25, "Next", fontsize=11, color=RED, fontweight="bold")
    ax.plot([10.05, 13.2], [2.12, 2.12], color=RED, linewidth=1.2)
    return save_fig(fig, "proposal_pipeline_diagram.png")


def make_sampling_panel() -> Path:
    patches = read_patch_rows()
    fig = plt.figure(figsize=(13, 5.1), facecolor=WHITE)
    gs = fig.add_gridspec(1, 2, width_ratios=[1.0, 1.25], wspace=0.12)
    ax = fig.add_subplot(gs[0, 0])
    ax.set_facecolor("#F0F4F8")

    seoul_box = (126.76, 37.41, 127.19, 37.70)
    ax.add_patch(
        mpatches.Rectangle(
            (seoul_box[0], seoul_box[1]),
            seoul_box[2] - seoul_box[0],
            seoul_box[3] - seoul_box[1],
            facecolor="#E5EDF6",
            edgecolor=NAVY,
            linewidth=1.2,
        )
    )

    for name, min_lon, min_lat, max_lon, max_lat in URBAN_POOLS:
        ax.add_patch(
            mpatches.Rectangle(
                (min_lon, min_lat),
                max_lon - min_lon,
                max_lat - min_lat,
                facecolor="none",
                edgecolor="#8FA6C6",
                linewidth=1.0,
                linestyle="--",
            )
        )
        ax.text((min_lon + max_lon) / 2, (min_lat + max_lat) / 2, name.split("_")[0], ha="center", va="center", fontsize=6, color="#52677F")

    handles = {}
    for row in patches:
        pid = row["patch_id"]
        split, color = SPLIT_INFO[pid]
        min_lon = float(row["min_lon"])
        min_lat = float(row["min_lat"])
        max_lon = float(row["max_lon"])
        max_lat = float(row["max_lat"])
        ax.add_patch(
            mpatches.Rectangle(
                (min_lon, min_lat),
                max_lon - min_lon,
                max_lat - min_lat,
                facecolor=color,
                edgecolor=color,
                linewidth=2.0,
                alpha=0.35,
            )
        )
        ax.scatter([(min_lon + max_lon) / 2], [(min_lat + max_lat) / 2], color=color, edgecolor=WHITE, s=45, zorder=5)
        ax.text((min_lon + max_lon) / 2, (min_lat + max_lat) / 2 + 0.01, patch_short_name(pid).split()[0], ha="center", va="bottom", fontsize=7.2, color=DARK, fontweight="bold")
        handles[split] = mpatches.Patch(facecolor=color, label=split)

    ax.set_xlim(126.75, 127.20)
    ax.set_ylim(37.40, 37.71)
    ax.set_title("9 urban pools and 6 selected patches", fontsize=12, color=NAVY, fontweight="bold")
    ax.set_xlabel("Longitude", fontsize=9, color=GRAY)
    ax.set_ylabel("Latitude", fontsize=9, color=GRAY)
    ax.grid(True, color="#D5DDE8", linewidth=0.6)
    ax.tick_params(labelsize=8, colors=GRAY)
    ax.legend(handles=[handles[k] for k in ("Train", "Val", "Test") if k in handles], loc="lower right", fontsize=8, frameon=True)

    outer = fig.add_subplot(gs[0, 1])
    outer.axis("off")
    sub = gs[0, 1].subgridspec(2, 3, wspace=0.05, hspace=0.15)
    for i, row in enumerate(patches):
        axp = fig.add_subplot(sub[i // 3, i % 3])
        pid = row["patch_id"]
        img_path = PATCH_TIF_DIR / f"{pid}.tif"
        img = read_patch_image(img_path)
        img = ImageOps.fit(img, (360, 250), method=Image.Resampling.LANCZOS)
        axp.imshow(img)
        axp.axis("off")
        split, color = SPLIT_INFO[pid]
        for sp in axp.spines.values():
            sp.set_visible(True)
            sp.set_color(color)
            sp.set_linewidth(3)
        axp.set_title(f"{patch_short_name(pid).title()}\n{split}", fontsize=8, color=color, fontweight="bold", pad=2)
    outer.set_title("Selected aerial patches", fontsize=12, color=NAVY, fontweight="bold", y=1.02)
    return save_fig(fig, "proposal_sampling_panel.png")


def make_split_map() -> Path:
    patches = read_patch_rows()
    fig, ax = plt.subplots(figsize=(6.2, 5), facecolor=WHITE)
    ax.set_facecolor("#F0F4F8")
    ax.add_patch(mpatches.Rectangle((126.76, 37.41), 0.43, 0.29, facecolor="#E5EDF6", edgecolor=NAVY, linewidth=1.2))
    handles = {}
    for row in patches:
        pid = row["patch_id"]
        split, color = SPLIT_INFO[pid]
        min_lon = float(row["min_lon"])
        min_lat = float(row["min_lat"])
        max_lon = float(row["max_lon"])
        max_lat = float(row["max_lat"])
        ax.add_patch(mpatches.Rectangle((min_lon, min_lat), max_lon - min_lon, max_lat - min_lat, facecolor=color, edgecolor=color, alpha=0.42, linewidth=2))
        ax.text((min_lon + max_lon) / 2, (min_lat + max_lat) / 2, patch_short_name(pid).title().replace(" ", "\n"), ha="center", va="center", fontsize=7, color=DARK, fontweight="bold")
        handles[split] = mpatches.Patch(facecolor=color, edgecolor=color, label=split)
    ax.set_xlim(126.75, 127.20)
    ax.set_ylim(37.40, 37.71)
    ax.grid(True, color="#D5DDE8", linewidth=0.6)
    ax.tick_params(labelsize=8, colors=GRAY)
    ax.set_xlabel("Longitude", fontsize=9, color=GRAY)
    ax.set_ylabel("Latitude", fontsize=9, color=GRAY)
    ax.set_title("Patch-level spatial split", fontsize=12, color=NAVY, fontweight="bold")
    ax.legend(handles=[handles[k] for k in ("Train", "Val", "Test") if k in handles], loc="lower right", fontsize=9)
    fig.tight_layout()
    return save_fig(fig, "proposal_split_map.png")


def make_tile_grid() -> Path:
    rows = read_csv_rows(SELECTED_CSV)
    fig, axes = plt.subplots(6, 6, figsize=(8.9, 8.7), facecolor=WHITE)
    for ax, row in zip(axes.flat, rows):
        file_stem = Path(row["file_name"]).stem
        rank = int(row["selection_rank"])
        tile_stem = f"rank{rank:02d}_{file_stem}"
        img_path = TILE_IMG_DIR / f"{tile_stem}.png"
        pid = row["patch_id"]
        split, color = SPLIT_INFO.get(pid, ("", MID_GRAY))
        img = open_rgb(img_path) if img_path.exists() else Image.new("RGB", (512, 512), "white")
        ax.imshow(img)
        ax.set_xticks([])
        ax.set_yticks([])
        for sp in ax.spines.values():
            sp.set_visible(True)
            sp.set_linewidth(2.2)
            sp.set_edgecolor(color)
        ax.text(
            0.03,
            0.94,
            tile_stem.split("_")[0].replace("rank", "#"),
            transform=ax.transAxes,
            color=WHITE,
            fontsize=8,
            fontweight="bold",
            bbox=dict(facecolor=color, alpha=0.95, edgecolor="none", pad=2),
        )
    for ax in axes.flat[len(rows) :]:
        ax.axis("off")
    handles = [mpatches.Patch(facecolor=color, label=f"{split}") for split, color in [("Train", BLUE), ("Val", ORANGE), ("Test", RED)]]
    fig.legend(handles=handles, loc="lower center", ncol=3, frameon=False, fontsize=10, bbox_to_anchor=(0.5, 0.01))
    fig.suptitle("36 labeling candidates (512 × 512 px) — color by split", fontsize=15, color=NAVY, fontweight="bold", y=0.995)
    fig.subplots_adjust(left=0.015, right=0.985, top=0.94, bottom=0.065, wspace=0.03, hspace=0.03)
    return save_fig(fig, "proposal_tile_grid.png")


def make_labeling_panel() -> Path:
    examples = [
        ("rank17_p04_seongdong_wangsimni_tile512_r00_c01", "Seongdong"),
        ("rank27_p02_songpa_jamsil_tile512_r03_c01", "Songpa"),
        ("rank35_p02_songpa_jamsil_tile512_r02_c02", "Songpa"),
    ]
    fig = plt.figure(figsize=(13, 5.7), facecolor=WHITE)
    gs = fig.add_gridspec(1, 2, width_ratios=[1.05, 1.35], wspace=0.12)

    ax_ui = fig.add_subplot(gs[0, 0])
    ax_ui.set_xlim(0, 1)
    ax_ui.set_ylim(0, 1)
    ax_ui.axis("off")
    ax_ui.add_patch(mpatches.FancyBboxPatch((0.02, 0.05), 0.96, 0.9, boxstyle="round,pad=0.015", facecolor="#F3F5F8", edgecolor=MID_GRAY))
    ax_ui.add_patch(mpatches.Rectangle((0.02, 0.86), 0.96, 0.09, facecolor=NAVY, edgecolor=NAVY))
    ax_ui.text(0.05, 0.905, "Tree canopy labeling app", fontsize=11, color=WHITE, fontweight="bold", va="center")
    ax_ui.add_patch(mpatches.FancyBboxPatch((0.06, 0.77), 0.22, 0.055, boxstyle="round,pad=0.01", facecolor=BLUE, edgecolor=BLUE))
    ax_ui.add_patch(mpatches.FancyBboxPatch((0.31, 0.77), 0.22, 0.055, boxstyle="round,pad=0.01", facecolor=ORANGE, edgecolor=ORANGE))
    ax_ui.text(0.17, 0.797, "Start polygon", ha="center", va="center", fontsize=8.5, color=WHITE, fontweight="bold")
    ax_ui.text(0.42, 0.797, "Save label", ha="center", va="center", fontsize=8.5, color=WHITE, fontweight="bold")

    stem = examples[0][0]
    _, img_path, mask_path = find_split_and_path(stem)
    img = open_rgb(img_path) if img_path else Image.new("RGB", (512, 512), "white")
    if mask_path:
        img = add_mask_overlay(img, Image.open(mask_path))
    img = ImageOps.fit(img, (420, 420), method=Image.Resampling.LANCZOS)
    ax_ui.imshow(img, extent=(0.09, 0.91, 0.12, 0.72), zorder=1)
    ax_ui.text(0.5, 0.08, "click/drag polygon vertices → orange canopy mask", ha="center", fontsize=8.5, color=GRAY)

    sub = gs[0, 1].subgridspec(2, 3, hspace=0.08, wspace=0.08)
    for col, (stem, title) in enumerate(examples):
        _, img_path, mask_path = find_split_and_path(stem)
        img = open_rgb(img_path) if img_path else Image.new("RGB", (512, 512), "white")
        mask = Image.open(mask_path) if mask_path else Image.new("L", img.size, 0)
        overlay = add_mask_overlay(img, mask)
        for row, panel_img in enumerate([img, overlay]):
            ax = fig.add_subplot(sub[row, col])
            ax.imshow(panel_img)
            ax.axis("off")
            if row == 0:
                ax.set_title(title, fontsize=10, color=NAVY, fontweight="bold", pad=3)
            if col == 0:
                ax.text(-0.05, 0.5, "Raw" if row == 0 else "Mask", rotation=90, va="center", ha="right", transform=ax.transAxes, fontsize=10, color=GRAY if row == 0 else ORANGE, fontweight="bold")
    fig.suptitle("Manual crown polygons → semantic canopy masks", fontsize=14, color=NAVY, fontweight="bold", y=0.99)
    return save_fig(fig, "proposal_labeling_panel.png")


def make_training_curve() -> Path:
    df = pd.read_csv(METRICS_CSV)
    best_idx = df["val_iou"].idxmax()
    best_ep = int(df.loc[best_idx, "epoch"])
    best_iou = float(df.loc[best_idx, "val_iou"])
    fig, ax = plt.subplots(figsize=(5.3, 3.1), facecolor=WHITE)
    ax.plot(df["epoch"], df["train_iou"], color=BLUE, linewidth=2, label="Train IoU")
    ax.plot(df["epoch"], df["val_iou"], color=NAVY, linewidth=2.2, label="Val IoU")
    ax.axvline(best_ep, color=RED, linewidth=1.2, linestyle="--")
    ax.scatter([best_ep], [best_iou], color=RED, s=35, zorder=3)
    ax.annotate(f"Best ep {best_ep}\nVal IoU {best_iou:.3f}", xy=(best_ep, best_iou), xytext=(best_ep + 4, best_iou + 0.035), fontsize=8, color=RED, arrowprops=dict(arrowstyle="->", color=RED, lw=0.8))
    ax.set_title("Fine-tuning curve", fontsize=11, color=NAVY, fontweight="bold")
    ax.set_xlabel("Epoch", fontsize=9, color=GRAY)
    ax.set_ylabel("IoU", fontsize=9, color=GRAY)
    ax.grid(True, color="#D7DCE3", linewidth=0.6)
    ax.tick_params(labelsize=8, colors=GRAY)
    ax.legend(fontsize=8, frameon=False)
    fig.tight_layout()
    return save_fig(fig, "proposal_training_curve.png")


def make_threshold_sweep() -> Path:
    df = pd.read_csv(THRESH_CSV)
    val = df[df["split"] == "val"].sort_values("threshold")
    test = df[df["split"] == "test"].sort_values("threshold")
    best_t = 0.55
    fig, ax = plt.subplots(figsize=(5.3, 3.1), facecolor=WHITE)
    ax.plot(val["threshold"], val["iou"], color=BLUE, linewidth=2, marker="o", ms=3.5, label="Val IoU")
    ax.plot(test["threshold"], test["iou"], color=NAVY, linewidth=2, marker="o", ms=3.5, label="Test IoU")
    ax.axvline(best_t, color=RED, linewidth=1.2, linestyle="--")
    test_iou = float(test.loc[np.isclose(test["threshold"], best_t), "iou"].iloc[0])
    ax.annotate(f"thr={best_t:.2f}\nTest IoU={test_iou:.3f}", xy=(best_t, test_iou), xytext=(best_t + 0.07, test_iou - 0.055), fontsize=8, color=RED, arrowprops=dict(arrowstyle="->", color=RED, lw=0.8))
    ax.set_title("Threshold sweep", fontsize=11, color=NAVY, fontweight="bold")
    ax.set_xlabel("Threshold", fontsize=9, color=GRAY)
    ax.set_ylabel("IoU", fontsize=9, color=GRAY)
    ax.grid(True, color="#D7DCE3", linewidth=0.6)
    ax.tick_params(labelsize=8, colors=GRAY)
    ax.legend(fontsize=8, frameon=False)
    fig.tight_layout()
    return save_fig(fig, "proposal_threshold_sweep.png")


def make_prediction_wide_panel() -> Path:
    """Create a landscape version of the long test preview sheet."""
    examples = [
        ("rank17_p04_seongdong_wangsimni_tile512_r00_c01", "Seongdong"),
        ("rank19_p04_seongdong_wangsimni_tile512_r00_c00", "Seongdong"),
        ("rank27_p02_songpa_jamsil_tile512_r03_c01", "Songpa"),
        ("rank35_p02_songpa_jamsil_tile512_r02_c02", "Songpa"),
    ]
    fig, axes = plt.subplots(3, len(examples), figsize=(13.4, 8.0), facecolor=WHITE)
    row_titles = [("Original", GRAY), ("Ground truth", GREEN), ("SegFormer pred", ORANGE)]
    for col, (stem, title) in enumerate(examples):
        _, img_path, mask_path = find_split_and_path(stem)
        pred_path = TCD_RUN_DIR / "threshold_pred_masks" / "test" / f"{stem}.png"
        img = open_rgb(img_path) if img_path else Image.new("RGB", (512, 512), "white")
        gt = Image.open(mask_path) if mask_path else Image.new("L", img.size, 0)
        pred = Image.open(pred_path) if pred_path.exists() else Image.new("L", img.size, 0)
        panels = [
            img,
            add_mask_overlay(img, gt, color=(47, 143, 104), alpha=0.45),
            add_mask_overlay(img, pred, color=(242, 142, 43), alpha=0.50),
        ]
        for row, panel in enumerate(panels):
            ax = axes[row, col]
            ax.imshow(panel)
            ax.axis("off")
            if row == 0:
                ax.set_title(title, fontsize=11, color=NAVY, fontweight="bold", pad=5)
            if col == 0:
                label, color = row_titles[row]
                ax.text(
                    -0.06,
                    0.5,
                    label,
                    rotation=90,
                    va="center",
                    ha="right",
                    transform=ax.transAxes,
                    fontsize=11,
                    color=color,
                    fontweight="bold",
                )
    fig.suptitle("Test prediction examples at threshold 0.55", fontsize=16, color=NAVY, fontweight="bold", y=0.985)
    fig.subplots_adjust(left=0.05, right=0.99, top=0.92, bottom=0.02, wspace=0.04, hspace=0.08)
    return save_fig(fig, "proposal_test_prediction_wide_panel.png")


def make_model_panel() -> Path:
    fig = plt.figure(figsize=(12.8, 4.8), facecolor=WHITE)
    gs = fig.add_gridspec(1, 2, width_ratios=[1.15, 1.0], wspace=0.15)
    ax = fig.add_subplot(gs[0, 0])
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 5)
    ax.axis("off")
    boxes = [
        (0.2, 2.0, 1.5, 1.0, "512×512\nRGB tile", LT_BLUE, BLUE),
        (2.2, 1.65, 1.7, 1.7, "MiT-B2\nhierarchical\nencoder", "#EADCF8", "#7030A0"),
        (4.5, 1.9, 1.6, 1.2, "MLP\nfusion\ndecoder", "#FFF2CC", "#9A6A00"),
        (6.7, 1.9, 1.5, 1.2, "Canopy\nlogits", LT_GREEN, GREEN),
        (8.7, 2.0, 1.1, 1.0, "Binary\nmask", LT_ORANGE, ORANGE),
    ]
    for x, y, w, h, text, fc, ec in boxes:
        ax.add_patch(mpatches.FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.04,rounding_size=0.05", facecolor=fc, edgecolor=ec, linewidth=1.6))
        ax.text(x + w / 2, y + h / 2, text, ha="center", va="center", fontsize=10, color=DARK, fontweight="bold")
    for x in [1.75, 4.0, 6.15, 8.25]:
        ax.annotate("", xy=(x + 0.35, 2.5), xytext=(x, 2.5), arrowprops=dict(arrowstyle="->", color=NAVY, lw=1.6))
    ax.text(0.2, 4.55, "SegFormer for aerial tree canopy segmentation", fontsize=13, color=NAVY, fontweight="bold")
    ax.text(0.2, 0.55, "Pretrained checkpoint: restor/tcd-segformer-mit-b2\nFine-tuned on Seoul V-World labels", fontsize=10.5, color=GRAY)

    sub = gs[0, 1].subgridspec(1, 2, width_ratios=[1, 1], wspace=0.07)
    sample = "rank17_p04_seongdong_wangsimni_tile512_r00_c01"
    _, img_path, mask_path = find_split_and_path(sample)
    img = open_rgb(img_path) if img_path else Image.new("RGB", (512, 512), "white")
    mask = Image.open(mask_path) if mask_path else Image.new("L", img.size, 0)
    overlay = add_mask_overlay(img, mask)
    for i, (panel, title) in enumerate([(img, "Input aerial tile"), (overlay, "Canopy mask target")]):
        axp = fig.add_subplot(sub[0, i])
        axp.imshow(panel)
        axp.axis("off")
        axp.set_title(title, fontsize=10, color=NAVY, fontweight="bold", pad=4)
    return save_fig(fig, "proposal_model_panel.png")


def make_research_design_schematic() -> Path:
    fig = plt.figure(figsize=(12.8, 4.8), facecolor=WHITE)
    gs = fig.add_gridspec(1, 2, width_ratios=[1.0, 1.0], wspace=0.16)
    ax1 = fig.add_subplot(gs[0, 0])
    ax1.set_facecolor("#F0F4F8")
    ax1.add_patch(mpatches.Rectangle((126.76, 37.41), 0.43, 0.29, facecolor="#E5EDF6", edgecolor=NAVY, linewidth=1.2))
    rng = np.random.default_rng(938)
    xs = rng.uniform(126.78, 127.17, 180)
    ys = rng.uniform(37.43, 37.68, 180)
    ax1.scatter(xs, ys, s=9, color=BLUE, alpha=0.55, edgecolor="none")
    ax1.set_xlim(126.75, 127.20)
    ax1.set_ylim(37.40, 37.71)
    ax1.grid(True, color="#D5DDE8", linewidth=0.6)
    ax1.tick_params(labelsize=8, colors=GRAY)
    ax1.set_title("S-DoT sensor distribution (schematic)", fontsize=12, color=NAVY, fontweight="bold")
    ax1.set_xlabel("Longitude", fontsize=9, color=GRAY)
    ax1.set_ylabel("Latitude", fontsize=9, color=GRAY)

    ax2 = fig.add_subplot(gs[0, 1])
    sample = "rank27_p02_songpa_jamsil_tile512_r03_c01"
    _, img_path, mask_path = find_split_and_path(sample)
    img = open_rgb(img_path) if img_path else Image.new("RGB", (512, 512), "white")
    mask = Image.open(mask_path) if mask_path else Image.new("L", img.size, 0)
    overlay = add_mask_overlay(img, mask, color=(47, 143, 104), alpha=0.42)
    ax2.imshow(overlay)
    ax2.axis("off")
    ax2.add_patch(mpatches.Circle((256, 256), 145, facecolor="none", edgecolor=ORANGE, linewidth=3))
    ax2.scatter([256], [256], s=80, color=RED, edgecolor=WHITE, linewidth=1.5, zorder=5)
    ax2.text(270, 246, "S-DoT\nsensor", fontsize=10, color=RED, fontweight="bold")
    ax2.text(310, 400, "50 m buffer\ncanopy overlap", fontsize=10, color=ORANGE, fontweight="bold")
    ax2.set_title("Canopy ratio joined to 50 m buffer", fontsize=12, color=NAVY, fontweight="bold")
    return save_fig(fig, "proposal_research_design_schematic.png")


def make_all_figures() -> dict[str, Path]:
    print("Generating proposal figures...")
    figs = {
        "motivation": make_motivation_diagram(),
        "pipeline": make_pipeline_diagram(),
        "sampling": make_sampling_panel(),
        "split_map": make_split_map(),
        "tile_grid": make_tile_grid(),
        "labeling": make_labeling_panel(),
        "training": make_training_curve(),
        "threshold": make_threshold_sweep(),
        "prediction_wide": make_prediction_wide_panel(),
        "model_panel": make_model_panel(),
        "research_design": make_research_design_schematic(),
    }
    return figs


def build_ppt(figs: dict[str, Path]) -> Path:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.util import Inches, Pt

    C_NAVY = RGBColor(0x0B, 0x1F, 0x3A)
    C_NAVY2 = RGBColor(0x00, 0x20, 0x60)
    C_BLUE = RGBColor(0x2E, 0x75, 0xB6)
    C_GREEN = RGBColor(0x2F, 0x8F, 0x68)
    C_ORANGE = RGBColor(0xF2, 0x8E, 0x2B)
    C_RED = RGBColor(0xC0, 0x00, 0x00)
    C_DARK = RGBColor(0x26, 0x26, 0x26)
    C_GRAY = RGBColor(0x66, 0x66, 0x66)
    C_LIGHT = RGBColor(0xEA, 0xEF, 0xF5)
    C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def add_text(
        slide,
        text: str,
        left: float,
        top: float,
        width: float,
        height: float,
        size: float = 12,
        color=C_DARK,
        bold: bool = False,
        italic: bool = False,
        align=PP_ALIGN.LEFT,
        font: str = "Lato",
    ):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        box.text_frame.word_wrap = True
        box.text_frame.margin_left = Inches(0.02)
        box.text_frame.margin_right = Inches(0.02)
        box.text_frame.margin_top = Inches(0.02)
        box.text_frame.margin_bottom = Inches(0.02)
        box.text_frame.vertical_anchor = MSO_ANCHOR.TOP
        box.text_frame.clear()
        for idx, line in enumerate(text.split("\n")):
            p = box.text_frame.paragraphs[0] if idx == 0 else box.text_frame.add_paragraph()
            p.alignment = align
            r = p.add_run()
            r.text = line
            r.font.name = font
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.italic = italic
            r.font.color.rgb = color
        return box

    def add_bullets(slide, lines: list[str], left: float, top: float, width: float, height: float, size: float = 12, color=C_DARK):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.02)
        tf.margin_right = Inches(0.02)
        tf.clear()
        for idx, line in enumerate(lines):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.level = 0
            p.space_after = Pt(5)
            run = p.add_run()
            run.text = line
            run.font.name = "Lato"
            run.font.size = Pt(size)
            run.font.color.rgb = color
        return box

    def add_rect(slide, left, top, width, height, fill, line=None, radius=False):
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        shp = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        if line:
            shp.line.color.rgb = line
            shp.line.width = Pt(0.8)
        else:
            shp.line.fill.background()
        return shp

    def add_line(slide, x1, y1, x2, y2, color=C_LIGHT, width=0.9):
        conn = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        conn.line.color.rgb = color
        conn.line.width = Pt(width)
        return conn

    def add_image(slide, path: Path, left: float, top: float, width: float, height: float | None = None):
        if height is None:
            slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))
            return
        im = Image.open(path)
        ratio = im.width / im.height
        box_ratio = width / height
        if ratio >= box_ratio:
            w = width
            h = width / ratio
            x = left
            y = top + (height - h) / 2
        else:
            h = height
            w = height * ratio
            x = left + (width - w) / 2
            y = top
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))

    def chrome(slide, title: str, number: int, section: str = "Preliminary Research Proposal"):
        add_rect(slide, 0, 0, 0.16, 7.5, C_NAVY)
        add_text(slide, section, 0.38, 0.23, 5.2, 0.3, size=9.5, color=C_NAVY2, bold=True)
        add_text(slide, str(number), 12.72, 0.18, 0.4, 0.3, size=10, color=C_GRAY, align=PP_ALIGN.RIGHT)
        add_text(slide, title, 0.38, 0.55, 11.9, 0.55, size=24, color=C_NAVY, bold=True)
        add_line(slide, 0.38, 1.16, 12.95, 1.16, C_LIGHT, 1.0)

    def add_simple_table(slide, rows, left, top, width, height, col_widths=None, font_size=9.5, header=True):
        tbl_shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(left), Inches(top), Inches(width), Inches(height))
        tbl = tbl_shape.table
        if col_widths:
            for idx, cw in enumerate(col_widths):
                tbl.columns[idx].width = Inches(cw)
        for r_idx, row in enumerate(rows):
            for c_idx, text in enumerate(row):
                cell = tbl.cell(r_idx, c_idx)
                cell.text = str(text)
                cell.margin_left = Inches(0.05)
                cell.margin_right = Inches(0.05)
                cell.margin_top = Inches(0.03)
                cell.margin_bottom = Inches(0.03)
                for p in cell.text_frame.paragraphs:
                    p.alignment = PP_ALIGN.CENTER if c_idx > 0 or r_idx == 0 else PP_ALIGN.LEFT
                    for run in p.runs:
                        run.font.name = "Lato"
                        run.font.size = Pt(font_size)
                        run.font.color.rgb = C_DARK
                        run.font.bold = bool(header and r_idx == 0)
                if header and r_idx == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = C_LIGHT
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = C_WHITE
        return tbl_shape

    # 0. Cover
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, 0.16, 7.5, C_NAVY)
    add_text(slide, "From Street-View to Aerial", 0.52, 1.25, 11.8, 0.75, size=34, color=C_NAVY, bold=True)
    add_text(
        slide,
        "Comparing Perceptual and Physical Urban Green/Built Metrics\nfor Microclimate Air Temperature Analysis",
        0.52,
        2.12,
        11.8,
        0.9,
        size=19,
        color=C_BLUE,
    )
    add_line(slide, 0.52, 3.42, 12.65, 3.42, C_LIGHT, 1.0)
    add_text(slide, "Preliminary Research Proposal", 0.52, 3.70, 8, 0.35, size=14, color=C_GRAY, bold=True)
    add_text(slide, "Seho Oh  |  Department of Urban Planning and Engineering  |  Yonsei University", 0.52, 4.30, 10.5, 0.35, size=12.5, color=C_DARK)
    add_text(slide, "Tree canopy map pilot: V-World aerial imagery, Seoul urban patches, SegFormer-B2 fine-tuning", 0.52, 5.0, 10.8, 0.35, size=11, color=C_GRAY, italic=True)
    add_text(slide, "2026", 11.8, 6.85, 0.8, 0.35, size=11, color=C_GRAY, align=PP_ALIGN.RIGHT)

    # 1. Motivation
    slide = prs.slides.add_slide(blank)
    chrome(slide, "Motivation", 1)
    add_image(slide, figs["motivation"], 0.38, 1.32, 12.55, 2.95)
    add_text(slide, "Prior work — Oh & Jung (2025)", 0.55, 4.45, 4.0, 0.3, size=12, color=C_BLUE, bold=True)
    add_bullets(
        slide,
        [
            "S-DoT AT × street-view metrics (GVI, BVI, SVF)",
            "LULC controlled → diurnal marginal effects estimated",
            "GVI cools in afternoon; BVI warms at night",
        ],
        0.55,
        4.86,
        5.2,
        1.35,
        size=10.6,
    )
    rows = [
        ["", "Oh & Jung (2025)", "This study"],
        ["Green", "GVI", "Aerial canopy map"],
        ["Built", "BVI", "Footprint & height"],
        ["Control", "LULC", "LULC"],
        ["Framework", "S-DoT + mixed-effects + Fourier", "Same"],
    ]
    add_simple_table(slide, rows, 6.0, 4.48, 6.75, 1.75, col_widths=[1.15, 2.8, 2.8], font_size=8.5)
    add_text(slide, "Do physical aerial measurements tell the same story?", 6.15, 6.42, 6.25, 0.32, size=12, color=C_RED, bold=True, italic=True)

    # 2. Prior work details
    slide = prs.slides.add_slide(blank)
    chrome(slide, "Prior Work — Street-View Urban Metrics", 2)
    add_text(slide, "Oh & Jung (2025): research setup", 0.55, 1.38, 5.7, 0.32, size=13, color=C_BLUE, bold=True)
    add_bullets(
        slide,
        [
            "Question: how local visual environment modifies near-surface air temperature around S-DoT sensors.",
            "Dependent variable: hourly S-DoT near-surface air temperature (AT/NAT).",
            "Core predictors: GVI, BVI, SVF derived from street-view semantic segmentation.",
            "LULC controlled so the street-level visual signal is not just a land-use proxy.",
            "Time structure modeled with Fourier terms to estimate diurnal marginal effects.",
            "Sensor-level random intercept absorbs persistent site-specific temperature differences.",
        ],
        0.55,
        1.85,
        5.9,
        3.35,
        size=10.7,
    )
    add_text(slide, "Metric interpretation", 6.85, 1.38, 5.7, 0.32, size=13, color=C_GREEN, bold=True)
    rows = [
        ["Metric", "Meaning", "Expected thermal pathway"],
        ["GVI", "Green View Index\nvisible vegetation share", "tree shade + evapotranspiration\nmainly afternoon cooling"],
        ["BVI", "Building View Index\nvisible building façade share", "urban enclosure + heat storage\nmainly nighttime warming"],
        ["SVF", "Sky View Factor\nvisible sky openness", "radiative cooling / solar exposure\ncontext-dependent"],
    ]
    add_simple_table(slide, rows, 6.85, 1.85, 5.85, 2.25, col_widths=[0.85, 2.35, 2.65], font_size=7.7)
    add_rect(slide, 6.85, 4.55, 5.85, 1.15, C_LIGHT, radius=True)
    add_text(slide, "Why this proposal extends it", 7.05, 4.72, 5.4, 0.25, size=11.5, color=C_NAVY, bold=True)
    add_text(
        slide,
        "Street-view metrics are perceptual: they describe what the sensor sees horizontally. This study asks whether physical aerial measurements around the same sensors reproduce or contradict those diurnal effects.",
        7.05,
        5.08,
        5.4,
        0.42,
        size=9.7,
        color=C_DARK,
    )
    add_text(slide, "Key prior finding: GVI cooling appears in the afternoon; BVI warming appears at night.", 0.75, 6.25, 11.8, 0.35, size=12, color=C_RED, bold=True, italic=True)

    # 3. Prior work controls
    slide = prs.slides.add_slide(blank)
    chrome(slide, "Prior Work — Data Controls & Weather Filtering", 3)
    add_text(slide, "Typical-summer-day weather controls", 0.55, 1.36, 5.7, 0.32, size=13, color=C_BLUE, bold=True)
    rows = [
        ["Filter", "Threshold", "Purpose"],
        ["Daily precipitation", "0.0 mm", "remove rainy-day cooling/wetting effects"],
        ["Mean cloud cover", "≤ 2/10", "keep clear-sky radiative conditions"],
        ["Mean wind speed", "≤ 2.0 m/s", "reduce advective mixing effects"],
        ["Daily max temperature", "< 35.0 °C", "exclude extreme heat days"],
    ]
    add_simple_table(slide, rows, 0.55, 1.9, 5.95, 1.85, col_widths=[1.75, 1.45, 2.75], font_size=8.1)
    add_text(slide, "Selected study days", 0.55, 4.25, 5.7, 0.32, size=13, color=C_GREEN, bold=True)
    add_bullets(
        slide,
        [
            "2023: Jun 24, Jul 01",
            "2024: Jun 13, Jun 26, Aug 31",
            "Unit: sensor × time-of-day window",
            "Time windows: dawn, morning, afternoon, night",
        ],
        0.55,
        4.72,
        5.7,
        1.55,
        size=10.8,
    )
    add_text(slide, "Modeling logic", 7.0, 1.36, 5.5, 0.32, size=13, color=C_BLUE, bold=True)
    add_bullets(
        slide,
        [
            "Use the same S-DoT sensor network and the same weather-screened thermal regime.",
            "Keep LULC controls so green/built effects are not merely area-composition effects.",
            "Use Fourier time terms to estimate when effects are strongest during the day.",
            "Compare perceptual SVI metrics with physical aerial metrics under the same framework.",
        ],
        7.0,
        1.85,
        5.6,
        1.8,
        size=10.8,
    )
    add_rect(slide, 7.0, 4.2, 5.6, 1.7, C_LIGHT, radius=True)
    add_text(slide, "Implication for this study", 7.2, 4.38, 5.2, 0.28, size=11.5, color=C_NAVY, bold=True)
    add_text(
        slide,
        "The aerial canopy model is only the measurement layer. The final causal/associational test should still happen in the same filtered S-DoT microclimate setting.",
        7.2,
        4.78,
        5.2,
        0.75,
        size=10.3,
        color=C_DARK,
    )

    # 4. Pipeline
    slide = prs.slides.add_slide(blank)
    chrome(slide, "Overall Pipeline", 4)
    add_image(slide, figs["pipeline"], 0.35, 1.4, 12.65, 2.65)
    add_text(slide, "Current output", 0.65, 4.45, 2.5, 0.3, size=13, color=C_BLUE, bold=True)
    add_bullets(
        slide,
        [
            "36 labeled tiles → binary semantic canopy masks",
            "restor/tcd-segformer-mit-b2 fine-tuned on Seoul aerial data",
            "Tree canopy prediction ready for 50 m sensor-buffer aggregation",
        ],
        0.65,
        4.88,
        5.65,
        1.35,
        size=11.5,
    )
    add_text(slide, "Next model stage", 7.0, 4.45, 2.5, 0.3, size=13, color=C_GREEN, bold=True)
    add_bullets(
        slide,
        [
            "DV: S-DoT hourly air temperature",
            "IV: canopy ratio + building footprint/height",
            "Controls: LULC + weather filters + temporal Fourier terms",
        ],
        7.0,
        4.88,
        5.5,
        1.35,
        size=11.5,
    )

    # 5. Study area
    slide = prs.slides.add_slide(blank)
    chrome(slide, "Study Area & Sampling", 5)
    add_image(slide, figs["sampling"], 0.35, 1.28, 8.1, 5.75)
    add_text(slide, "Sampling strategy", 8.75, 1.35, 4.0, 0.3, size=13, color=C_BLUE, bold=True)
    add_bullets(
        slide,
        [
            "9 predefined Seoul urban pools",
            "Hand-picked to avoid mountains, rivers, highways",
            "Random seed 20260419 → 6 patches selected",
            "Exclusion mask: UPIS park zones + urban nature park zones",
            "277 features union/dissolved → single GeoJSON",
            "≥5% bbox overlap → resample",
        ],
        8.75,
        1.78,
        4.15,
        2.95,
        size=10.5,
    )
    add_text(slide, "Data limitation", 8.75, 5.00, 4.0, 0.3, size=12.5, color=C_RED, bold=True)
    add_text(
        slide,
        "V-World WMTS JPEG has no shooting date or sensor metadata.\nFuture work: NGII national orthoimagery with acquisition metadata.",
        8.75,
        5.42,
        4.1,
        0.92,
        size=10.5,
        color=C_DARK,
    )

    # 6. Tile selection
    slide = prs.slides.add_slide(blank)
    chrome(slide, "Tile Selection", 6)
    add_image(slide, figs["tile_grid"], 0.34, 1.22, 6.75, 6.0)
    add_text(slide, "125 total tiles → 36 labeling candidates", 8.55, 1.36, 4.3, 0.3, size=13, color=C_BLUE, bold=True)
    rows = [["Filter", "Removed"], ["Edge tiles", "25"], ["Geometry overlap", "1"], ["Forest-like (HSV)", "0"], ["Per-patch cap (max 8)", "remainder"]]
    add_simple_table(slide, rows, 8.55, 1.86, 4.25, 1.55, col_widths=[2.7, 1.45], font_size=9)
    add_text(slide, "Why this matters", 8.55, 3.78, 4.3, 0.3, size=12.5, color=C_GREEN, bold=True)
    add_bullets(
        slide,
        [
            "Keeps labeling effort urban-focused",
            "Reduces park/mountain leakage",
            "Maintains patch-level spatial split",
            "Small pilot set, but high-quality masks",
        ],
        8.55,
        4.2,
        4.25,
        1.55,
        size=10.8,
    )

    # 7. Labeling
    slide = prs.slides.add_slide(blank)
    chrome(slide, "Labeling", 7)
    add_image(slide, figs["labeling"], 0.35, 1.22, 8.75, 5.95)
    add_text(slide, "Manual labeling", 9.0, 1.35, 3.8, 0.3, size=13, color=C_BLUE, bold=True)
    add_bullets(
        slide,
        [
            "Crown boundary polygons",
            "Class: tree only",
            "Excluded: shrubs, grass, shadows",
            "475 canopy polygons",
            "36 tiles, 2 empty tiles",
        ],
        9.0,
        1.78,
        3.8,
        1.75,
        size=10.8,
    )
    add_text(slide, "Strategy pivot", 9.0, 4.02, 3.8, 0.3, size=12.5, color=C_ORANGE, bold=True)
    add_text(
        slide,
        "Initially attempted instance segmentation.\nOverlapping crowns were too hard to separate consistently, so labels were unioned/rasterized into semantic masks.",
        9.0,
        4.45,
        3.85,
        1.25,
        size=10.6,
        color=C_DARK,
    )
    add_text(slide, "Known limitation: dense crown boundaries may be inconsistent.", 9.0, 6.15, 3.85, 0.3, size=9.8, color=C_GRAY, italic=True)

    # 8. Dataset split
    slide = prs.slides.add_slide(blank)
    chrome(slide, "Dataset Split", 8)
    add_image(slide, figs["split_map"], 0.55, 1.35, 5.3, 5.4)
    add_text(slide, "Patch-level spatial split", 6.15, 1.38, 5.9, 0.3, size=13, color=C_BLUE, bold=True)
    rows = [
        ["Split", "Patches", "Tiles", "Canopy ratio"],
        ["Train", "Yeongdeungpo, Mapo, Gangnam", "24", "4.3%"],
        ["Val", "Guro Digital", "5", "3.4%"],
        ["Test", "Songpa, Seongdong", "7", "9.0%"],
    ]
    add_simple_table(slide, rows, 6.15, 1.9, 6.55, 1.8, col_widths=[1.0, 3.55, 0.8, 1.2], font_size=8.5)
    add_bullets(
        slide,
        [
            "Source patch is the split unit, not individual tile.",
            "Prevents spatial leakage between adjacent tiles.",
            "Test set has higher canopy ratio, so recall/precision balance matters.",
        ],
        6.15,
        4.15,
        6.2,
        1.25,
        size=11.3,
    )

    # 9. Model
    slide = prs.slides.add_slide(blank)
    chrome(slide, "Model", 9)
    add_image(slide, figs["model_panel"], 0.35, 1.28, 8.1, 5.7)
    add_text(slide, "restor/tcd-segformer-mit-b2", 8.45, 1.35, 4.35, 0.3, size=13, color=C_BLUE, bold=True)
    add_bullets(
        slide,
        [
            "Pretrained on OpenAerialMap Tree Canopy Dataset",
            "Designed for aerial tree canopy segmentation",
            "Pretrain resolution: 10 cm/px",
            "Current V-World data: ~25 cm/px",
        ],
        8.45,
        1.78,
        4.3,
        1.65,
        size=10.7,
    )
    add_text(slide, "Fine-tuning setup", 8.45, 3.85, 4.35, 0.3, size=12.5, color=C_GREEN, bold=True)
    add_bullets(
        slide,
        [
            "512 px / batch 4 / 50 epochs",
            "AdamW lr 5e-5",
            "Canopy class weight 5",
            "Best epoch: 15",
            "Best val IoU: 0.412",
        ],
        8.45,
        4.28,
        4.3,
        1.75,
        size=10.7,
    )

    # 10. Results
    slide = prs.slides.add_slide(blank)
    chrome(slide, "Results — Metrics", 10)
    add_image(slide, figs["training"], 0.35, 1.28, 3.2, 2.05)
    add_image(slide, figs["threshold"], 3.76, 1.28, 3.2, 2.05)
    add_image(slide, PREVIEW_SHEET, 10.1, 1.26, 2.7, 5.85)
    rows = [
        ["Model", "Setting", "IoU", "Dice", "Precision", "Recall"],
        ["SegFormer-B0", "fine-tuned", "0.507", "0.673", "0.679", "0.666"],
        ["TCD-B2", "zero-shot argmax", "0.625", "0.769", "0.853", "0.700"],
        ["TCD-B2", "zero-shot thr=0.30", "0.654", "0.791", "0.763", "0.821"],
        ["TCD-B2", "fine-tuned argmax", "0.686", "0.814", "0.748", "0.892"],
        ["TCD-B2", "fine-tuned thr=0.55", "0.691", "0.817", "0.763", "0.879"],
    ]
    add_simple_table(slide, rows, 0.35, 3.72, 9.35, 2.1, col_widths=[1.55, 2.8, 1.0, 1.0, 1.45, 1.25], font_size=8.0)
    add_text(slide, "36 labeled tiles → IoU +0.184 over scratch baseline", 0.55, 6.30, 7.8, 0.32, size=12.2, color=C_RED, bold=True)
    add_text(slide, "Preview: original / ground truth / SegFormer prediction", 10.0, 6.93, 2.9, 0.22, size=7.5, color=C_GRAY, align=PP_ALIGN.CENTER)

    # 11. Results preview
    slide = prs.slides.add_slide(blank)
    chrome(slide, "Results — Test Prediction Examples", 11)
    add_image(slide, figs["prediction_wide"], 0.35, 1.25, 12.55, 5.85)
    add_text(
        slide,
        "Landscape panel reconstructed from the vertical preview: original tile, manual canopy mask, and thresholded SegFormer prediction.",
        0.55,
        6.98,
        11.8,
        0.25,
        size=9.5,
        color=C_GRAY,
        italic=True,
    )

    # 12. Next research design
    slide = prs.slides.add_slide(blank)
    chrome(slide, "Next: Research Design", 12)
    add_image(slide, figs["research_design"], 0.35, 1.28, 7.4, 3.2)
    rows = [
        ["", "Oh & Jung (2025)", "This study"],
        ["Green metric", "GVI", "Canopy ratio"],
        ["Built metric", "BVI", "Footprint & height"],
        ["Control", "LULC + weather filters", "LULC + same filters"],
        ["Sensors", "938 S-DoT", "938 S-DoT"],
        ["Model", "Mixed-effects + Fourier", "Mixed-effects + Fourier"],
    ]
    add_simple_table(slide, rows, 8.0, 1.35, 4.95, 2.15, col_widths=[1.35, 1.8, 1.8], font_size=7.8)
    add_text(slide, "Key questions", 0.55, 4.82, 2.3, 0.3, size=13, color=C_BLUE, bold=True)
    add_bullets(
        slide,
        [
            "1. Does aerial canopy ratio show the same afternoon cooling pattern as GVI?",
            "2. Does building footprint & height show the same nighttime warming as BVI?",
            "3. Where do the two representations diverge, and what does that reveal?",
        ],
        0.55,
        5.25,
        7.1,
        1.4,
        size=10.9,
    )
    add_text(slide, "Future extension", 8.1, 4.25, 2.8, 0.3, size=12.5, color=C_GREEN, bold=True)
    add_text(
        slide,
        "DSM-based shadow casting → diurnal trade-off:\ntree shade at noon vs building shade in morning/evening.",
        8.1,
        4.68,
        4.65,
        0.85,
        size=10.8,
        color=C_DARK,
    )
    add_rect(slide, 8.1, 5.93, 4.55, 0.62, C_LIGHT, radius=True)
    add_text(slide, "Next deliverable: citywide canopy ratio around each S-DoT sensor", 8.28, 6.08, 4.2, 0.25, size=10, color=C_NAVY, bold=True)

    prs.save(OUT_PPT)
    print(f"\nSaved PPT: {OUT_PPT.relative_to(ROOT)}")
    return OUT_PPT


def main() -> None:
    figs = make_all_figures()
    print("\nBuilding PowerPoint...")
    build_ppt(figs)


if __name__ == "__main__":
    main()
