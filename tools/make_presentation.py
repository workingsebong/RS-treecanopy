"""
Generate presentation PPT matching template style.
Run: conda activate svi_segformer && python tools/make_presentation.py
"""

import json
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.patches as patches
import numpy as np
import pandas as pd
from PIL import Image

ROOT    = Path(__file__).parent.parent
OUT_FIG = ROOT / "outputs" / "figures" / "presentation"  # intermediate figures
OUT_PPT = ROOT / "papers" / "presentation"
OUT_FIG.mkdir(parents=True, exist_ok=True)
OUT_PPT.mkdir(parents=True, exist_ok=True)

# ── color palette (matches template) ──────────────────────────────────────
NAVY    = "#142D52"
NAVY2   = "#002060"
WHITE   = "#FFFFFF"
LGRAY   = "#F5F5F5"
MGRAY   = "#E0E0E0"
DKGRAY  = "#444444"
ACCENT  = "#2E75B6"   # medium blue for highlights
RED_HL  = "#C00000"

# ── data paths ────────────────────────────────────────────────────────────
METRICS_CSV   = ROOT / "outputs/segformer/tree_canopy_semantic_tcd_mit_b2_finetune_e50/metrics.csv"
THRESH_CSV    = ROOT / "outputs/segformer/tree_canopy_semantic_tcd_mit_b2_finetune_e50/threshold_metrics.csv"
DATASET_META  = ROOT / "data/raw/random_seoul_urban/dataset_metadata.json"
TILE_DIR      = ROOT / "data/processed/labeling_candidates_512_random_seoul_urban/images"
MASK_DIR      = ROOT / "data/processed/semantic_seg_512_random_seoul_urban/masks"
IMG_DIR       = ROOT / "data/processed/semantic_seg_512_random_seoul_urban/images"
PREVIEW_SHEET = ROOT / "outputs/segformer/tree_canopy_semantic_tcd_mit_b2_finetune_e50/test_threshold_preview_sheet.jpg"
ZEROSHOT_SH   = ROOT / "outputs/segformer/tree_canopy_semantic_tcd_mit_b2_zeroshot/test_argmax_preview_sheet.jpg"
B0_SHEET      = ROOT / "outputs/segformer/tree_canopy_semantic_mit_b0_gpu_e50/test_threshold_preview_sheet.jpg"

SPLIT_INFO = {
    "p01_guro_digital":         ("Val",   "#FF8C00"),
    "p02_songpa_jamsil":        ("Test",  "#C00000"),
    "p03_yeongdeungpo_mullae":  ("Train", "#2E75B6"),
    "p04_seongdong_wangsimni":  ("Test",  "#C00000"),
    "p05_mapo_hongdae_yeonnam": ("Train", "#2E75B6"),
    "p06_gangnam_teheran":      ("Train", "#2E75B6"),
}

def fig_style(fig, ax_list=None):
    fig.patch.set_facecolor(WHITE)
    if ax_list:
        for ax in ax_list:
            ax.set_facecolor(WHITE)
            ax.tick_params(colors=DKGRAY, labelsize=9)
            for sp in ax.spines.values():
                sp.set_edgecolor(MGRAY)
            ax.xaxis.label.set_color(DKGRAY)
            ax.yaxis.label.set_color(DKGRAY)
            ax.grid(True, color=MGRAY, linewidth=0.5, alpha=0.8)


# ── Figure 1: Training curve ───────────────────────────────────────────────
def make_training_curve():
    df = pd.read_csv(METRICS_CSV)
    best_ep = int(df["val_iou"].idxmax()) + 1
    best_iou = df["val_iou"].max()

    fig, ax = plt.subplots(figsize=(6, 3.2), facecolor=WHITE)
    ax.plot(df["epoch"], df["train_iou"], color=ACCENT,  lw=1.8, label="Train IoU")
    ax.plot(df["epoch"], df["val_iou"],   color=NAVY,    lw=2.0, label="Val IoU")
    ax.axvline(best_ep, color=RED_HL, lw=1.2, ls="--")
    ax.annotate(f"Best ep {best_ep}\nVal IoU {best_iou:.3f}",
                xy=(best_ep, best_iou),
                xytext=(best_ep + 3, best_iou - 0.06),
                color=RED_HL, fontsize=8,
                arrowprops=dict(arrowstyle="->", color=RED_HL, lw=0.8))
    ax.set_xlabel("Epoch", fontsize=9, color=DKGRAY)
    ax.set_ylabel("IoU", fontsize=9, color=DKGRAY)
    ax.set_title("Fine-tuning: Val IoU per Epoch", fontsize=10, color=NAVY, fontweight="bold")
    fig_style(fig, [ax])
    ax.legend(fontsize=8, framealpha=0.9, edgecolor=MGRAY)
    fig.tight_layout()
    p = OUT_FIG / "training_curve.png"
    fig.savefig(p, dpi=150, bbox_inches="tight", facecolor=WHITE)
    plt.close(fig)
    print(f"  {p.name}")
    return p


# ── Figure 2: Threshold sweep ──────────────────────────────────────────────
def make_threshold_sweep():
    df   = pd.read_csv(THRESH_CSV)
    val  = df[df["split"] == "val"].sort_values("threshold")
    test = df[df["split"] == "test"].sort_values("threshold")
    best_t = 0.55

    fig, ax = plt.subplots(figsize=(6, 3.2), facecolor=WHITE)
    ax.plot(val["threshold"],  val["iou"],  color=ACCENT, lw=1.8, label="Val IoU",  marker="o", ms=3)
    ax.plot(test["threshold"], test["iou"], color=NAVY,   lw=1.8, label="Test IoU", marker="o", ms=3)
    ax.axvline(best_t, color=RED_HL, lw=1.2, ls="--")
    t_iou = test.loc[abs(test["threshold"] - best_t) < 0.01, "iou"].values
    if len(t_iou):
        ax.annotate(f"thr={best_t}\nTest IoU={t_iou[0]:.3f}",
                    xy=(best_t, t_iou[0]),
                    xytext=(best_t + 0.06, t_iou[0] - 0.04),
                    color=RED_HL, fontsize=8,
                    arrowprops=dict(arrowstyle="->", color=RED_HL, lw=0.8))
    ax.set_xlabel("Threshold", fontsize=9, color=DKGRAY)
    ax.set_ylabel("IoU",       fontsize=9, color=DKGRAY)
    ax.set_title("Threshold Sweep (val-selected → test)", fontsize=10, color=NAVY, fontweight="bold")
    fig_style(fig, [ax])
    ax.legend(fontsize=8, framealpha=0.9, edgecolor=MGRAY)
    fig.tight_layout()
    p = OUT_FIG / "threshold_sweep.png"
    fig.savefig(p, dpi=150, bbox_inches="tight", facecolor=WHITE)
    plt.close(fig)
    print(f"  {p.name}")
    return p


# ── Figure 3: Patch location map ──────────────────────────────────────────
def make_patch_map():
    with open(DATASET_META) as f:
        meta = json.load(f)

    fig, ax = plt.subplots(figsize=(4.5, 5), facecolor=WHITE)
    ax.set_facecolor("#EEF2F7")

    # Seoul rough outline
    seoul_lon = [126.764, 126.764, 127.183, 127.183, 126.764]
    seoul_lat = [37.413, 37.701, 37.701, 37.413, 37.413]
    ax.fill(seoul_lon, seoul_lat, color="#D8E4F0", zorder=1)
    ax.plot(seoul_lon, seoul_lat, color=NAVY, lw=1.2, zorder=2)

    handles = {}
    short_names = {
        "p01_guro_digital":         "Guro\nDigital",
        "p02_songpa_jamsil":        "Songpa\nJamsil",
        "p03_yeongdeungpo_mullae":  "Yeong-\ndeungpo",
        "p04_seongdong_wangsimni":  "Seong-\ndong",
        "p05_mapo_hongdae_yeonnam": "Mapo\nHongdae",
        "p06_gangnam_teheran":      "Gangnam\nTeheran",
    }
    for p in meta["patches"]:
        pid   = p["patch_id"]
        label, color = SPLIT_INFO[pid]
        cx = (p["min_lon"] + p["max_lon"]) / 2
        cy = (p["min_lat"] + p["max_lat"]) / 2
        w  = p["max_lon"] - p["min_lon"]
        h  = p["max_lat"] - p["min_lat"]
        rect = mpatches.FancyBboxPatch(
            (p["min_lon"], p["min_lat"]), w, h,
            boxstyle="square,pad=0.0001",
            linewidth=2, edgecolor=color,
            facecolor=color + "33", zorder=3
        )
        ax.add_patch(rect)
        ax.text(cx, cy, short_names[pid], ha="center", va="center",
                fontsize=6.5, color=color, zorder=4, fontweight="bold",
                linespacing=1.2)
        if label not in handles:
            handles[label] = mpatches.Patch(color=color, label=label)

    ax.set_xlim(126.75, 127.22)
    ax.set_ylim(37.40, 37.72)
    ax.set_xlabel("Longitude", fontsize=8, color=DKGRAY)
    ax.set_ylabel("Latitude",  fontsize=8, color=DKGRAY)
    ax.tick_params(colors=DKGRAY, labelsize=7)
    for sp in ax.spines.values(): sp.set_edgecolor(MGRAY)
    ax.set_title("6 Urban Patches — Seoul", fontsize=10, color=NAVY, fontweight="bold", pad=6)
    ax.legend(handles=list(handles.values()), fontsize=8, framealpha=0.9,
              edgecolor=MGRAY, loc="lower right")
    ax.grid(True, color=MGRAY, linewidth=0.4, alpha=0.5)
    fig.tight_layout()
    p = OUT_FIG / "patch_map.png"
    fig.savefig(p, dpi=150, bbox_inches="tight", facecolor=WHITE)
    plt.close(fig)
    print(f"  {p.name}")
    return p


# ── Figure 4: Representative tile sample (3 per split = 9 tiles) ──────────
def make_tile_sample():
    """Show 3 tiles per split, bigger and readable."""
    picks = {
        "Train": [
            "rank01_p03_yeongdeungpo_mullae_tile512_r03_c03",
            "rank08_p05_mapo_hongdae_yeonnam_tile512_r00_c00",
            "rank26_p06_gangnam_teheran_tile512_r00_c00",
        ],
        "Val": [
            "rank20_p01_guro_digital_tile512_r03_c03",
            "rank25_p01_guro_digital_tile512_r02_c00",
            "rank33_p01_guro_digital_tile512_r01_c00",
        ],
        "Test": [
            "rank17_p04_seongdong_wangsimni_tile512_r00_c01",
            "rank27_p02_songpa_jamsil_tile512_r03_c01",
            "rank35_p02_songpa_jamsil_tile512_r02_c02",
        ],
    }
    split_colors = {"Train": "#2E75B6", "Val": "#FF8C00", "Test": "#C00000"}

    fig, axes = plt.subplots(3, 3, figsize=(7.5, 7.5), facecolor=WHITE)
    for row, (split, names) in enumerate(picks.items()):
        col_color = split_colors[split]
        for col, name in enumerate(names):
            ax = axes[row, col]
            img_path = TILE_DIR / f"{name}.png"
            if img_path.exists():
                ax.imshow(Image.open(img_path).convert("RGB"))
            ax.axis("off")
            # colored border
            for sp in ax.spines.values():
                sp.set_visible(True)
                sp.set_edgecolor(col_color)
                sp.set_linewidth(3)
            if col == 0:
                ax.set_ylabel(split, fontsize=11, color=col_color,
                              fontweight="bold", labelpad=6)
                ax.yaxis.set_visible(True)
                ax.yaxis.set_ticks([])

    fig.suptitle("Labeling Candidates (9 of 36) — Color by Split",
                 fontsize=11, color=NAVY, fontweight="bold", y=1.01)
    fig.tight_layout(pad=0.4)
    p = OUT_FIG / "tile_sample.png"
    fig.savefig(p, dpi=150, bbox_inches="tight", facecolor=WHITE)
    plt.close(fig)
    print(f"  {p.name}")
    return p


# ── Figure 5: Label overlay (image | mask overlay), 3 examples ────────────
def make_label_overlay():
    examples = [
        ("train", "rank01_p03_yeongdeungpo_mullae_tile512_r03_c03",   "Yeongdeungpo"),
        ("train", "rank08_p05_mapo_hongdae_yeonnam_tile512_r00_c00",  "Mapo / Hongdae"),
        ("test",  "rank17_p04_seongdong_wangsimni_tile512_r00_c01",   "Seongdong"),
    ]

    fig, axes = plt.subplots(2, 3, figsize=(9, 6), facecolor=WHITE)
    for col, (split, name, title) in enumerate(examples):
        img_path  = IMG_DIR  / split / f"{name}.png"
        mask_path = MASK_DIR / split / f"{name}.png"

        img  = Image.open(img_path).convert("RGB")  if img_path.exists()  else None
        mask = Image.open(mask_path).convert("L")   if mask_path.exists() else None

        # top row: raw image
        ax_top = axes[0, col]
        if img: ax_top.imshow(img)
        ax_top.axis("off")
        ax_top.set_title(title, fontsize=10, color=NAVY, fontweight="bold", pad=4)

        # bottom row: image + canopy overlay
        ax_bot = axes[1, col]
        if img:
            ax_bot.imshow(img)
            if mask:
                m = np.array(mask)
                overlay = np.zeros((*m.shape, 4), dtype=np.uint8)
                overlay[m > 0] = [46, 139, 87, 180]   # sea-green, semi-transparent
                ax_bot.imshow(overlay)
        ax_bot.axis("off")

    axes[0, 0].set_ylabel("Aerial image", fontsize=9, color=DKGRAY,
                           labelpad=4, visible=True)
    axes[0, 0].yaxis.set_ticks([])
    axes[1, 0].set_ylabel("+ Canopy mask", fontsize=9, color="#2E8B57",
                           labelpad=4, visible=True)
    axes[1, 0].yaxis.set_ticks([])

    fig.suptitle("Labeled Tile Examples  (green = canopy)",
                 fontsize=11, color=NAVY, fontweight="bold", y=1.01)
    fig.tight_layout(pad=0.5)
    p = OUT_FIG / "label_overlay.png"
    fig.savefig(p, dpi=150, bbox_inches="tight", facecolor=WHITE)
    plt.close(fig)
    print(f"  {p.name}")
    return p


# ── Figure 6: Model comparison bar chart ──────────────────────────────────
def make_model_comparison():
    models = ["B0\nscratch", "TCD-B2\nzero-shot\nargmax",
              "TCD-B2\nzero-shot\nthr=0.30",
              "TCD-B2\nfine-tune\nargmax",
              "TCD-B2\nfine-tune\nthr=0.55"]
    iou  = [0.507, 0.625, 0.654, 0.686, 0.691]
    dice = [0.673, 0.769, 0.791, 0.814, 0.817]
    prec = [0.679, 0.853, 0.763, 0.748, 0.763]
    rec  = [0.666, 0.700, 0.821, 0.892, 0.879]

    x = np.arange(len(models))
    w = 0.19
    colors = ["#2E75B6", "#70AD47", "#FFC000", "#ED7D31"]

    fig, ax = plt.subplots(figsize=(8.5, 3.8), facecolor=WHITE)
    b1 = ax.bar(x - 1.5*w, iou,  w, label="IoU",       color=colors[0], zorder=3)
    b2 = ax.bar(x - 0.5*w, dice, w, label="Dice",      color=colors[1], zorder=3)
    b3 = ax.bar(x + 0.5*w, prec, w, label="Precision", color=colors[2], zorder=3)
    b4 = ax.bar(x + 1.5*w, rec,  w, label="Recall",    color=colors[3], zorder=3)

    # highlight best model
    for bar in [b1[-1], b2[-1]]:
        bar.set_edgecolor(RED_HL)
        bar.set_linewidth(2)

    # value labels on IoU bars
    for i, (bar, v) in enumerate(zip(b1, iou)):
        ax.text(bar.get_x() + bar.get_width()/2, v + 0.008, f"{v:.3f}",
                ha="center", va="bottom", fontsize=7,
                color=RED_HL if i == len(iou)-1 else DKGRAY,
                fontweight="bold" if i == len(iou)-1 else "normal")

    ax.set_xticks(x)
    ax.set_xticklabels(models, fontsize=8, color=DKGRAY)
    ax.set_ylim(0, 1.05)
    ax.set_ylabel("Score", fontsize=9, color=DKGRAY)
    ax.set_title("Test Set Performance Comparison", fontsize=11, color=NAVY, fontweight="bold")
    fig_style(fig, [ax])
    ax.legend(fontsize=8, framealpha=0.9, edgecolor=MGRAY, ncol=4,
              loc="upper left")
    fig.tight_layout()
    p = OUT_FIG / "model_comparison.png"
    fig.savefig(p, dpi=150, bbox_inches="tight", facecolor=WHITE)
    plt.close(fig)
    print(f"  {p.name}")
    return p


# ── Figure 7: Pipeline diagram ────────────────────────────────────────────
def make_pipeline_diagram():
    steps = [
        ("V-World WMTS\nAerial Imagery\n~0.25 m/px", "#D8E4F0", NAVY),
        ("Urban Patch\nSampling\n(6 patches)", "#D8E4F0", NAVY),
        ("512 px Tile Split\n→ 36 Candidates", "#D8E4F0", NAVY),
        ("Manual\nLabeling\n475 polygons", "#D8E4F0", NAVY),
        ("Binary Semantic\nMask", "#D8E4F0", NAVY),
        ("SegFormer\nFine-tuning", "#C6EFCE", "#375623"),
        ("Tree Canopy\nMap", "#FFE699", "#7F6000"),
        ("S-DoT AT\nMixed-Effects\nModel  [Next]", "#FCE4D6", "#843C0C"),
    ]

    fig, ax = plt.subplots(figsize=(13, 2.4), facecolor=WHITE)
    ax.set_xlim(0, len(steps) * 1.6)
    ax.set_ylim(0, 2)
    ax.axis("off")
    fig.patch.set_facecolor(WHITE)

    box_w, box_h = 1.3, 1.4
    gap = 0.3

    for i, (text, fc, tc) in enumerate(steps):
        x0 = i * (box_w + gap) + gap / 2
        y0 = 0.3
        rect = patches.FancyBboxPatch(
            (x0, y0), box_w, box_h,
            boxstyle="round,pad=0.05",
            facecolor=fc, edgecolor=NAVY, linewidth=1.2
        )
        ax.add_patch(rect)
        ax.text(x0 + box_w/2, y0 + box_h/2, text,
                ha="center", va="center", fontsize=8, color=tc,
                fontweight="bold", linespacing=1.4)
        if i < len(steps) - 1:
            ax.annotate("", xy=(x0 + box_w + gap, y0 + box_h/2),
                        xytext=(x0 + box_w, y0 + box_h/2),
                        arrowprops=dict(arrowstyle="->", color=NAVY, lw=1.5))

    fig.tight_layout(pad=0.2)
    p = OUT_FIG / "pipeline_diagram.png"
    fig.savefig(p, dpi=150, bbox_inches="tight", facecolor=WHITE)
    plt.close(fig)
    print(f"  {p.name}")
    return p


# ── Build PPT ────────────────────────────────────────────────────────────────
def build_ppt(fig_paths: dict):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Emu

    C_NAVY  = RGBColor(0x14, 0x2D, 0x52)
    C_NAVY2 = RGBColor(0x00, 0x20, 0x60)
    C_DARK  = RGBColor(0x26, 0x26, 0x26)
    C_GRAY  = RGBColor(0x59, 0x59, 0x59)
    C_LGRAY = RGBColor(0xE0, 0xE0, 0xE0)
    C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    C_BLUE  = RGBColor(0x2E, 0x75, 0xB6)
    C_RED   = RGBColor(0xC0, 0x00, 0x00)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def new_slide():
        return prs.slides.add_slide(blank)

    def add_rect(sl, left, top, width, height, color):
        shp = sl.shapes.add_shape(1, left, top, width, height)
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.line.fill.background()
        return shp

    def add_line(sl, x1, y, x2, color=C_LGRAY, width_pt=0.75):
        from pptx.util import Pt
        conn = sl.shapes.add_connector(1, x1, y, x2, y)
        conn.line.color.rgb = color
        conn.line.width = Pt(width_pt)
        return conn

    def txb(sl, text, left, top, width, height,
            size=12, bold=False, color=C_DARK,
            align=PP_ALIGN.LEFT, italic=False):
        box = sl.shapes.add_textbox(left, top, width, height)
        box.word_wrap = True
        tf = box.text_frame
        tf.word_wrap = True
        for i, line in enumerate(text.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            r = p.add_run()
            r.text = line
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.italic = italic
            r.font.color.rgb = color
        return box

    def add_img(sl, path, left, top, width=None, height=None):
        if not Path(path).exists():
            return
        kw = {}
        if width:  kw["width"]  = width
        if height: kw["height"] = height
        sl.shapes.add_picture(str(path), left, top, **kw)

    def page_chrome(sl, title, part_label="RS · Microclimate"):
        """Left navy bar + title + divider + part label + slide number placeholder."""
        # left bar
        add_rect(sl, Inches(0), Inches(0.53), Inches(0.52), Inches(6.97), C_NAVY)
        # title
        txb(sl, title, Inches(0.65), Inches(0.58), Inches(10.5), Inches(0.6),
            size=26, bold=True, color=C_NAVY)
        # horizontal rule
        add_line(sl, Inches(0.65), Inches(1.22), Inches(12.85),
                 color=C_LGRAY, width_pt=1.0)
        # part label top-left above bar
        txb(sl, part_label, Inches(0.55), Inches(0.18), Inches(5), Inches(0.35),
            size=11, bold=True, color=C_NAVY2)

    # ── SLIDE 1 — Title ───────────────────────────────────────────────────
    sl = new_slide()
    add_rect(sl, Inches(0), Inches(0), Inches(0.52), Inches(7.5), C_NAVY)
    add_line(sl, Inches(0.65), Inches(3.8), Inches(12.85), C_LGRAY, 1.0)
    txb(sl, "From Street-View to Aerial",
        Inches(0.65), Inches(1.5), Inches(12.2), Inches(1.0),
        size=36, bold=True, color=C_NAVY)
    txb(sl, "Comparing Perceptual and Physical Green/Built Metrics\nfor Microclimate Air Temperature Analysis",
        Inches(0.65), Inches(2.45), Inches(12.2), Inches(1.0),
        size=18, bold=False, color=C_BLUE)
    txb(sl, "Preliminary Research Proposal",
        Inches(0.65), Inches(3.95), Inches(8), Inches(0.45),
        size=14, color=C_GRAY)
    txb(sl, "Seho Oh   |   Dept. of Urban Planning & Engineering, Yonsei University",
        Inches(0.65), Inches(4.45), Inches(10), Inches(0.45),
        size=13, bold=True, color=C_DARK)
    txb(sl, "CRP6588 Urban Heat Analysis  |  2026-04-22",
        Inches(0.65), Inches(4.95), Inches(8), Inches(0.4),
        size=12, color=C_GRAY)

    # ── SLIDE 2 — Motivation ─────────────────────────────────────────────
    sl = new_slide()
    page_chrome(sl, "Motivation", "Microclimate · Urban Form")
    # left panel
    txb(sl, "Prior work — Oh & Jung (2025)",
        Inches(0.65), Inches(1.35), Inches(5.5), Inches(0.4),
        size=13, bold=True, color=C_BLUE)
    txb(sl,
        "· S-DoT AT × street-view GVI, BLD, SVF\n"
        "· LULC controlled (12 classes)\n"
        "· Mixed-effects + Fourier time terms\n"
        "· Sensor random intercept\n\n"
        "Key findings:\n"
        "· GVI → cooling effect in afternoon\n"
        "· BLD → warming effect at night",
        Inches(0.65), Inches(1.8), Inches(5.5), Inches(2.5),
        size=12, color=C_DARK)
    txb(sl, "BUT",
        Inches(0.65), Inches(4.35), Inches(1), Inches(0.45),
        size=16, bold=True, color=C_RED)
    txb(sl,
        "GVI is perceptual — horizontal visibility of trees from street level\n"
        "Does actual overhead canopy coverage explain AT the same way?",
        Inches(0.65), Inches(4.85), Inches(5.5), Inches(0.9),
        size=12, color=C_DARK, italic=True)
    # vertical divider
    add_rect(sl, Inches(6.4), Inches(1.3), Inches(0.04), Inches(5.7), C_LGRAY)
    # right panel — comparison table
    txb(sl, "This study",
        Inches(6.6), Inches(1.35), Inches(6.5), Inches(0.4),
        size=13, bold=True, color=C_NAVY)
    rows = [
        ("", "Oh & Jung (2025)", "This study"),
        ("Green",   "GVI (SVI)",            "Aerial canopy map"),
        ("Built",   "BLD (SVI)",            "Footprint & height"),
        ("Control", "LULC",                 "LULC"),
        ("Sensor",  "938 S-DoT",            "938 S-DoT"),
        ("Model",   "ME + Fourier",         "ME + Fourier"),
    ]
    row_h = 0.52
    for ri, (a, b, c) in enumerate(rows):
        y = Inches(1.85 + ri * row_h)
        if ri == 0:
            txb(sl, b, Inches(7.5), y, Inches(2.5), Inches(row_h), size=10, bold=True, color=C_NAVY)
            txb(sl, c, Inches(10.2), y, Inches(2.9), Inches(row_h), size=10, bold=True, color=C_NAVY)
        else:
            txb(sl, a, Inches(6.6), y, Inches(0.9), Inches(row_h), size=11, bold=True, color=C_GRAY)
            txb(sl, b, Inches(7.5), y, Inches(2.5), Inches(row_h), size=11, color=C_DARK)
            txb(sl, c, Inches(10.2), y, Inches(2.9), Inches(row_h), size=11, bold=True, color=C_BLUE)
        if ri > 0:
            add_line(sl, Inches(6.6), y, Inches(13.1), C_LGRAY, 0.5)
    txb(sl, "→  Same framework, different spatial representation of green/built",
        Inches(6.6), Inches(5.2), Inches(6.5), Inches(0.45),
        size=11, bold=True, color=C_NAVY, italic=True)

    # ── SLIDE 3 — Pipeline ────────────────────────────────────────────────
    sl = new_slide()
    page_chrome(sl, "Overall Pipeline", "RS · Microclimate")
    add_img(sl, fig_paths["pipeline"],
            Inches(0.6), Inches(1.35), width=Inches(12.5))
    txb(sl,
        "Current status: Canopy map complete (IoU 0.691)  ·  S-DoT AT preprocessed  ·  Mixed-effects model pending",
        Inches(0.65), Inches(6.5), Inches(12.2), Inches(0.5),
        size=11, color=C_GRAY, italic=True)

    # ── SLIDE 4 — Study Area ─────────────────────────────────────────────
    sl = new_slide()
    page_chrome(sl, "Study Area & Sampling", "RS · Microclimate")
    add_img(sl, fig_paths["patch_map"],
            Inches(0.6), Inches(1.35), height=Inches(5.8))
    add_rect(sl, Inches(5.6), Inches(1.3), Inches(0.04), Inches(5.7), C_LGRAY)
    txb(sl, "Sampling strategy",
        Inches(5.8), Inches(1.35), Inches(7.2), Inches(0.4),
        size=13, bold=True, color=C_BLUE)
    txb(sl,
        "· 9 predefined Seoul urban pool zones\n"
        "  (avoids mountains, rivers, highways)\n"
        "· Random seed 20260419  →  6 patches\n\n"
        "Exclusion mask\n"
        "· UPIS park zones (UQ153, 208 features)\n"
        "· Urban nature park zones (UQ142, 69)\n"
        "· Union/dissolved → 277 features\n"
        "· ≥5% bbox overlap → resample\n\n"
        "⚠  Data limitation\n"
        "V-World WMTS JPEG:\nno shooting date or sensor metadata\n\n"
        "Future: NGII orthoimagery (12/25 cm)\nwith full acquisition metadata",
        Inches(5.8), Inches(1.85), Inches(7.2), Inches(5.2),
        size=11.5, color=C_DARK)

    # ── SLIDE 5 — Tile Selection ──────────────────────────────────────────
    sl = new_slide()
    page_chrome(sl, "Tile Selection — 36 Labeling Candidates", "RS · Microclimate")
    add_img(sl, fig_paths["tile_sample"],
            Inches(0.6), Inches(1.35), height=Inches(5.7))
    add_rect(sl, Inches(8.2), Inches(1.3), Inches(0.04), Inches(5.7), C_LGRAY)
    txb(sl, "125 total tiles  →  36 candidates",
        Inches(8.4), Inches(1.35), Inches(4.7), Inches(0.4),
        size=13, bold=True, color=C_BLUE)
    txb(sl,
        "Filters applied\n"
        "· Edge tiles removed:  25\n"
        "· Geometry overlap:     1\n"
        "· Forest-like (HSV):    0\n"
        "· Per-patch cap (≤8)\n\n"
        "Spatial split (patch-level)\n"
        "  Train   p03 p05 p06   24 tiles   4.3%\n"
        "  Val     p01            5 tiles   3.4%\n"
        "  Test    p02 p04        7 tiles   9.0%\n"
        "                     (canopy pixel ratio)\n\n"
        "Split unit = source patch\n"
        "→ prevents spatial leakage\n"
        "   between adjacent tiles",
        Inches(8.4), Inches(1.85), Inches(4.7), Inches(5.2),
        size=11.5, color=C_DARK)

    # ── SLIDE 6 — Labeling ────────────────────────────────────────────────
    sl = new_slide()
    page_chrome(sl, "Manual Labeling", "RS · Microclimate")
    add_img(sl, fig_paths["label_overlay"],
            Inches(0.6), Inches(1.35), width=Inches(9.0))
    add_rect(sl, Inches(9.8), Inches(1.3), Inches(0.04), Inches(5.7), C_LGRAY)
    txb(sl, "Labeling protocol",
        Inches(10.0), Inches(1.35), Inches(3.1), Inches(0.4),
        size=13, bold=True, color=C_BLUE)
    txb(sl,
        "· Crown boundary polygons\n"
        "· Class: tree only\n"
        "· Excluded: shrubs,\n"
        "  grass, shadows\n\n"
        "Result\n"
        "475 canopy polygons\n"
        "36 tiles  (2 empty)\n\n"
        "⚠  Strategy pivot\n"
        "Started with instance seg\n"
        "→ overlapping crowns\n"
        "   too hard to separate\n"
        "→ switched to semantic\n"
        "→ polygons union/\n"
        "   rasterized to mask\n"
        "→ possible boundary\n"
        "   inconsistency in\n"
        "   dense areas",
        Inches(10.0), Inches(1.85), Inches(3.1), Inches(5.2),
        size=11, color=C_DARK)

    # ── SLIDE 7 — Model ───────────────────────────────────────────────────
    sl = new_slide()
    page_chrome(sl, "Model — restor/tcd-segformer-mit-b2", "RS · Microclimate")
    txb(sl, "Why this model?",
        Inches(0.65), Inches(1.35), Inches(6.0), Inches(0.4),
        size=13, bold=True, color=C_BLUE)
    txb(sl,
        "· Pretrained on OpenAerialMap Tree Canopy Dataset\n"
        "  (ETH Zurich / Restor — global high-res aerial imagery)\n"
        "· Designed specifically for aerial tree canopy segmentation\n"
        "· B0–B5 backbone variants\n"
        "· Pretrain resolution: 10 cm/px\n"
        "  Current data: 25 cm/px   ⚠ scale mismatch\n\n"
        "Fine-tuning setup\n"
        "· 512 px  ·  batch 4  ·  50 epochs  ·  AdamW lr 5e-5\n"
        "· Canopy class weight 5  (pixel ratio ~5%)\n"
        "· Best epoch: 15   Best val IoU: 0.412",
        Inches(0.65), Inches(1.85), Inches(6.0), Inches(4.5),
        size=12, color=C_DARK)
    add_img(sl, fig_paths["training_curve"],
            Inches(6.8), Inches(1.3), width=Inches(6.2))

    # ── SLIDE 8 — Results ─────────────────────────────────────────────────
    sl = new_slide()
    page_chrome(sl, "Results", "RS · Microclimate")
    add_img(sl, fig_paths["model_comparison"],
            Inches(0.6), Inches(1.3), width=Inches(7.8))
    add_img(sl, fig_paths["threshold_sweep"],
            Inches(8.55), Inches(1.3), width=Inches(4.55))
    add_img(sl, PREVIEW_SHEET,
            Inches(0.6), Inches(4.1), width=Inches(12.5))

    # ── SLIDE 9 — Preview comparison ─────────────────────────────────────
    sl = new_slide()
    page_chrome(sl, "Prediction Preview — Test Set", "RS · Microclimate")
    txb(sl, "SegFormer-B0 (scratch)  IoU 0.507",
        Inches(0.65), Inches(1.3), Inches(4.1), Inches(0.35),
        size=10, bold=True, color=C_GRAY)
    add_img(sl, B0_SHEET, Inches(0.65), Inches(1.65), width=Inches(3.9))

    txb(sl, "TCD-B2  zero-shot  IoU 0.625",
        Inches(4.75), Inches(1.3), Inches(4.1), Inches(0.35),
        size=10, bold=True, color=C_BLUE)
    add_img(sl, ZEROSHOT_SH, Inches(4.75), Inches(1.65), width=Inches(3.9))

    txb(sl, "TCD-B2  fine-tuned  thr=0.55  IoU 0.691  ✓",
        Inches(8.85), Inches(1.3), Inches(4.3), Inches(0.35),
        size=10, bold=True, color=C_NAVY)
    add_img(sl, PREVIEW_SHEET, Inches(8.85), Inches(1.65), width=Inches(4.2))

    # ── SLIDE 10 — Next ───────────────────────────────────────────────────
    sl = new_slide()
    page_chrome(sl, "Next Steps — Research Design", "RS · Microclimate")
    txb(sl, "Canopy map  →  AT analysis",
        Inches(0.65), Inches(1.35), Inches(6.5), Inches(0.4),
        size=13, bold=True, color=C_BLUE)
    txb(sl,
        "DV        S-DoT hourly AT  (938 sensors, Seoul summer 2023–2024)\n"
        "IV         Tree canopy ratio within 50 m buffer\n"
        "Control  LULC (12 classes)  +  Building coverage ratio\n"
        "Model   Linear mixed-effects + Fourier time terms\n"
        "              + GVI×Fourier interaction  (sensor random intercept)",
        Inches(0.65), Inches(1.85), Inches(6.5), Inches(2.2),
        size=12, color=C_DARK)
    txb(sl, "Key questions",
        Inches(0.65), Inches(4.15), Inches(6.5), Inches(0.4),
        size=13, bold=True, color=C_BLUE)
    txb(sl,
        "1.  Does aerial canopy ratio show the same afternoon cooling as GVI?\n"
        "2.  Does building footprint & height show the same nighttime warming as BLD?\n"
        "3.  Where do the two representations diverge?",
        Inches(0.65), Inches(4.65), Inches(6.5), Inches(1.8),
        size=12, color=C_DARK)
    add_rect(sl, Inches(7.3), Inches(1.3), Inches(0.04), Inches(5.7), C_LGRAY)
    txb(sl, "Framework comparison",
        Inches(7.5), Inches(1.35), Inches(5.6), Inches(0.4),
        size=13, bold=True, color=C_BLUE)
    rows2 = [
        ("",        "Oh & Jung (2025)", "This study"),
        ("Green",   "GVI (SVI)",            "Aerial canopy map"),
        ("Built",   "BLD (SVI)",            "Footprint & height"),
        ("Control", "LULC",                 "LULC"),
        ("Sensors", "938 S-DoT",            "938 S-DoT"),
        ("Model",   "ME + Fourier",         "ME + Fourier"),
    ]
    for ri, (a, b, c) in enumerate(rows2):
        y = Inches(1.85 + ri * 0.52)
        if ri == 0:
            txb(sl, b, Inches(8.5), y, Inches(2.5), Inches(0.5), size=10, bold=True, color=C_NAVY)
            txb(sl, c, Inches(11.1), y, Inches(2.0), Inches(0.5), size=10, bold=True, color=C_NAVY)
        else:
            txb(sl, a, Inches(7.5), y, Inches(1.0), Inches(0.5), size=11, bold=True, color=C_GRAY)
            txb(sl, b, Inches(8.5), y, Inches(2.5), Inches(0.5), size=11, color=C_DARK)
            txb(sl, c, Inches(11.1), y, Inches(2.0), Inches(0.5), size=11, bold=True, color=C_BLUE)
        if ri > 0:
            add_line(sl, Inches(7.5), y, Inches(13.1), C_LGRAY, 0.5)
    txb(sl, "Future: DSM-based shadow casting\n→ diurnal trade-off: tree shade (noon) vs building shade (morning/evening)",
        Inches(7.5), Inches(5.2), Inches(5.6), Inches(1.1),
        size=11, color=C_NAVY, bold=True, italic=True)

    out = ROOT / "papers" / "presentation" / "tree_canopy_presentation.pptx"
    prs.save(str(out))
    print(f"\n✓  Saved: {out}")
    return out


# ── Main ──────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    print("Generating figures...")
    figs = {}
    figs["training_curve"]  = make_training_curve()
    figs["threshold_sweep"] = make_threshold_sweep()
    figs["patch_map"]       = make_patch_map()
    figs["tile_sample"]     = make_tile_sample()
    figs["label_overlay"]   = make_label_overlay()
    figs["model_comparison"]= make_model_comparison()
    figs["pipeline"]        = make_pipeline_diagram()

    print("\nBuilding PPT...")
    build_ppt(figs)
